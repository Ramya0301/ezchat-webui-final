import requests
import logging
import ftfy
import pandas as pd
from pptx import Presentation
import sys

from langchain_community.document_loaders import (
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    UnstructuredEPubLoader,
    UnstructuredExcelLoader,
    UnstructuredMarkdownLoader,
    UnstructuredPowerPointLoader,
    UnstructuredRSTLoader,
    UnstructuredXMLLoader,
    YoutubeLoader,
)
from langchain_core.documents import Document
from open_webui.env import SRC_LOG_LEVELS, GLOBAL_LOG_LEVEL

logging.basicConfig(stream=sys.stdout, level=GLOBAL_LOG_LEVEL)
log = logging.getLogger(__name__)
log.setLevel(SRC_LOG_LEVELS["RAG"])

known_source_ext = [
    "go",
    "py",
    "java",
    "sh",
    "bat",
    "ps1",
    "cmd",
    "js",
    "ts",
    "css",
    "cpp",
    "hpp",
    "h",
    "c",
    "cs",
    "sql",
    "log",
    "ini",
    "pl",
    "pm",
    "r",
    "dart",
    "dockerfile",
    "env",
    "php",
    "hs",
    "hsc",
    "lua",
    "nginxconf",
    "conf",
    "m",
    "mm",
    "plsql",
    "perl",
    "rb",
    "rs",
    "db2",
    "scala",
    "bash",
    "swift",
    "vue",
    "svelte",
    "msg",
    "ex",
    "exs",
    "erl",
    "tsx",
    "jsx",
    "hs",
    "lhs",
]


class TikaLoader:
    def __init__(self, url, file_path, mime_type=None):
        self.url = url
        self.file_path = file_path
        self.mime_type = mime_type

    def load(self) -> list[Document]:
        with open(self.file_path, "rb") as f:
            data = f.read()

        if self.mime_type is not None:
            headers = {"Content-Type": self.mime_type}
        else:
            headers = {}

        endpoint = self.url
        if not endpoint.endswith("/"):
            endpoint += "/"
        endpoint += "tika/text"

        r = requests.put(endpoint, data=data, headers=headers)

        if r.ok:
            raw_metadata = r.json()
            text = raw_metadata.get("X-TIKA:content", "<No text content found>")

            if "Content-Type" in raw_metadata:
                headers["Content-Type"] = raw_metadata["Content-Type"]

            log.debug("Tika extracted text: %s", text)

            return [Document(page_content=text, metadata=headers)]
        else:
            raise Exception(f"Error calling Tika: {r.reason}")

class CustomExcelLoader:
    def __init__(self, file_path: str, chunk_size: int = 1000, max_rows: int = None):
        self.file_path = file_path
        self.chunk_size = chunk_size
        self.max_rows = max_rows

    def load(self) -> list[Document]:
        try:
            excel_file = pd.ExcelFile(self.file_path)
            documents = []

            for sheet_name in excel_file.sheet_names:
                df = self._load_sheet(sheet_name)
                if df is not None:  # Skip empty or invalid sheets
                    content = self._convert_to_string(sheet_name, df)
                    metadata = self._generate_metadata(sheet_name, df)
                    documents.append(Document(page_content=content, metadata=metadata))
            
            return documents
        except Exception as e:
            logging.error(f"Error loading Excel file '{self.file_path}': {str(e)}")
            raise

    def _load_sheet(self, sheet_name: str) -> pd.DataFrame:
        try:
            df = pd.read_excel(
                self.file_path,
                sheet_name=sheet_name,
                nrows=self.max_rows,
                usecols=lambda col: not col.startswith("Unnamed"),
            )
            if df.empty or len(df.columns) == 0:
                logging.info(f"Skipping empty sheet: {sheet_name}")
                return None
            return df
        except Exception as e:
            logging.warning(f"Error processing sheet '{sheet_name}': {str(e)}")
            return None

    def _convert_to_string(self, sheet_name: str, df: pd.DataFrame) -> str:
        content = f"Sheet: {sheet_name}\n"
        content += df.head(self.chunk_size).to_string(index=False)
        return content

    def _generate_metadata(self, sheet_name: str, df: pd.DataFrame) -> dict:
        metadata = {
            "source": self.file_path,
            "sheet_name": sheet_name,
        }
        try:
            metadata.update({
                "row_count": len(df),
                "column_count": len(df.columns),
            })
        except Exception as e:
            logging.warning(f"Could not generate metadata for sheet '{sheet_name}': {str(e)}")
        return metadata
    
class CustomPowerPointLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        try:
            prs = Presentation(self.file_path)
            text_content = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():  # Only add non-empty text
                            slide_text.append(shape.text.strip())
                
                if slide_text:  # Only add slides with text
                    text_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(text_content)
            return [Document(page_content=full_text, metadata={"source": self.file_path})]
        except Exception as e:
            log.error(f"Error loading PowerPoint file: {e}")
            raise e

class Loader:
    def __init__(self, engine: str = "", **kwargs):
        self.engine = engine
        self.kwargs = kwargs

    def load(
        self, filename: str, file_content_type: str, file_path: str
    ) -> list[Document]:
        loader = self._get_loader(filename, file_content_type, file_path)
        docs = loader.load()

        return [
            Document(
                page_content=ftfy.fix_text(doc.page_content), metadata=doc.metadata
            )
            for doc in docs
        ]

    def _get_loader(self, filename: str, file_content_type: str, file_path: str):
        file_ext = filename.split(".")[-1].lower()

        if self.engine == "tika" and self.kwargs.get("TIKA_SERVER_URL"):
            if file_ext in known_source_ext or (
                file_content_type and file_content_type.find("text/") >= 0
            ):
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                loader = TikaLoader(
                    url=self.kwargs.get("TIKA_SERVER_URL"),
                    file_path=file_path,
                    mime_type=file_content_type,
                )
        else:
            if file_ext == "pdf":
                loader = PyPDFLoader(
                    file_path, extract_images=self.kwargs.get("PDF_EXTRACT_IMAGES")
                )
            elif file_ext == "csv":
                loader = CSVLoader(file_path)
            elif file_ext == "rst":
                loader = UnstructuredRSTLoader(file_path, mode="elements")
            elif file_ext == "xml":
                loader = UnstructuredXMLLoader(file_path)
            elif file_ext in ["htm", "html"]:
                loader = BSHTMLLoader(file_path, open_encoding="unicode_escape")
            elif file_ext == "md":
                loader = TextLoader(file_path, autodetect_encoding=True)
            elif file_content_type == "application/epub+zip":
                loader = UnstructuredEPubLoader(file_path)
            elif (
                file_content_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                or file_ext == "docx"
            ):
                loader = Docx2txtLoader(file_path)
            elif file_content_type in [
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ] or file_ext in ["xls", "xlsx"]:
                loader = CustomExcelLoader(file_path)
            elif file_content_type in [
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ] or file_ext in ["ppt", "pptx"]:
                loader = CustomPowerPointLoader(file_path)
            elif file_ext == "msg":
                loader = OutlookMessageLoader(file_path)
            elif file_ext in known_source_ext or (
                file_content_type and file_content_type.find("text/") >= 0
            ):
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                loader = TextLoader(file_path, autodetect_encoding=True)

        return loader
